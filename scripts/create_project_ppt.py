from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/presentation/tinyhear_project_intro.pptx")


BG = RGBColor(248, 250, 252)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
LINE = RGBColor(203, 213, 225)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size=24, bold=False, color=INK):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text="", size=24, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.35, 12.2, 0.45, title, 28, True, INK)
    if subtitle:
        add_textbox(slide, 0.58, 0.84, 11.9, 0.32, subtitle, 11, False, MUTED)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.color.rgb = LINE


def add_footer(slide, idx):
    add_textbox(slide, 11.85, 7.02, 0.9, 0.25, f"{idx:02d}", 10, False, MUTED, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, x + 0.2, y + 0.14, w - 0.35, 0.28, title, 15, True, INK)
    tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.48), Inches(w - 0.35), Inches(h - 0.58))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = body if isinstance(body, list) else [body]
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(11.5)
        p.font.color.rgb = MUTED
        p.space_after = Pt(4)
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=15, color=INK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
    return tb


def add_metric(slide, x, y, w, label, value, note, color=BLUE):
    add_textbox(slide, x, y, w, 0.22, label, 10, False, MUTED, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 0.26, w, 0.48, value, 25, True, color, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + 0.78, w, 0.24, note, 9.5, False, MUTED, PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=10.5):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(239, 246, 255) if r == 0 else WHITE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = INK if r == 0 else MUTED
    return table


def add_flow(slide, y, labels, colors):
    x = 0.7
    width = 1.55
    gap = 0.23
    for i, label in enumerate(labels):
        shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(width), Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.color.rgb = colors[i]
        add_textbox(slide, x + 0.05, y + 0.18, width - 0.1, 0.22, label, 10.5, True, WHITE, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_textbox(slide, x + width + 0.04, y + 0.2, 0.15, 0.18, ">", 14, True, MUTED, PP_ALIGN.CENTER)
        x += width + gap


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(12):
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = BG
        slides.append(s)

    # 1
    s = slides[0]
    add_textbox(s, 0.75, 0.9, 11.9, 0.55, "微听双麦降噪", 34, True, INK)
    add_textbox(s, 0.78, 1.55, 10.5, 0.35, "TinyHear Dual-Mic Denoising", 18, False, BLUE)
    add_textbox(s, 0.8, 2.25, 10.9, 0.55, "面向助听器 / 耳戴设备上行链路的端侧 AI 降噪原型", 23, True, INK)
    add_bullets(
        s,
        0.86,
        3.12,
        7.8,
        1.7,
        [
            "双麦输入，单通道增强输出",
            "16 kHz，256 点 FFT，64 samples hop，4 ms 帧移",
            "当前主线模型 137,984 参数，面向实时和后续 INT8/C 部署",
        ],
        16,
        MUTED,
    )
    add_card(s, 9.0, 3.0, 3.25, 1.35, "当前推荐版本", ["teacher_deepfilter_conservative_on_wind_eval", "checkpoint: runs/.../best.pt"], GREEN)
    add_footer(s, 1)

    # 2
    s = slides[1]
    add_title(s, "项目问题与约束", "为什么这个任务不是普通离线语音增强")
    add_card(s, 0.75, 1.55, 3.55, 2.0, "应用痛点", ["环境噪声、风噪、气流声会影响上行通话清晰度", "用户主观感受集中在底噪和“呼呼声”"], RED)
    add_card(s, 4.9, 1.55, 3.55, 2.0, "端侧限制", ["算力、SRAM、Flash、功耗有限", "需要低延迟、因果或近似因果实时处理"], AMBER)
    add_card(s, 9.05, 1.55, 3.55, 2.0, "双麦特点", ["近距双麦 ILD 较弱", "IPD / coherence / 空间统计更有价值"], BLUE)
    add_bullets(s, 1.0, 4.45, 11.2, 1.1, ["阶段目标：在小模型约束下，验证双麦空间前端 + 神经网络增强是否能形成可展示、可部署的完整闭环。"], 18, INK)
    add_footer(s, 2)

    # 3
    s = slides[2]
    add_title(s, "当前系统流程", "实时推理链路和网页试听闭环")
    add_flow(s, 1.65, ["双麦 waveform", "STFT", "空间特征", "MWF 前端", "Tiny DeepFilter", "iSTFT 输出", "网页 Demo"], [BLUE, CYAN, GREEN, AMBER, BLUE, CYAN, GREEN])
    add_card(s, 0.8, 3.0, 5.8, 2.1, "推理时实际运行", ["读取双通道 mixture", "提取 IPD / coherence 等空间特征", "coherence_mwf 生成参考通道", "TinyDeepFilterTCN 输出 gain 和 deep-filter 系数"], BLUE)
    add_card(s, 6.95, 3.0, 5.55, 2.1, "展示闭环", ["离线评估输出 metrics.json", "实时增强脚本生成 wav", "网页中直接对比 Noisy / Enhanced / Clean"], GREEN)
    add_footer(s, 3)

    # 4
    s = slides[3]
    add_title(s, "模型结构", "TinyDeepFilterTCN：小模型 + 复数域/多帧滤波")
    rows = [
        ["项目", "配置"],
        ["模型", "TinyDeepFilterTCN"],
        ["参数量", "137,984"],
        ["channels / blocks", "96 / 8"],
        ["FFT / hop", "256 / 64 samples"],
        ["DeepFilter", "df_bins=64, df_order=3"],
        ["spatial frontend", "coherence_mwf"],
    ]
    add_table(s, 0.85, 1.55, 4.6, 3.8, rows, [1.8, 2.8], 11)
    add_card(s, 6.0, 1.55, 6.1, 1.55, "输出 1：band gain", ["用于整体谱幅度抑制，保持轻量、稳定，适合端侧实现。"], BLUE)
    add_card(s, 6.0, 3.45, 6.1, 1.55, "输出 2：deep-filter 复数系数", ["对低频区域做多帧复数滤波，缓解单帧幅度 mask 带来的相位问题和残留风噪。"], CYAN)
    add_footer(s, 4)

    # 5
    s = slides[4]
    add_title(s, "数据与评估集", "从通用噪声到风噪场景")
    rows = [
        ["阶段", "Clean speech", "Noise", "作用"],
        ["Baseline", "合成语音", "合成噪声", "打通流程"],
        ["YESNO + DEMAND", "YESNO", "DEMAND", "小规模验证"],
        ["ARCTIC + DEMAND", "CMU ARCTIC", "DEMAND", "主基线"],
        ["当前", "CMU ARCTIC", "DEMAND + Zenodo Wind", "风噪优化主线"],
    ]
    add_table(s, 0.75, 1.45, 11.8, 2.45, rows, [2.3, 2.4, 2.8, 4.3], 10.3)
    add_card(s, 1.0, 4.55, 5.2, 1.25, "训练目录", ["data/arctic_wind_demand_flat/"], BLUE)
    add_card(s, 7.0, 4.55, 5.2, 1.25, "固定评估集", ["data/arctic_wind_eval/val", "data/arctic_demand_eval/val"], GREEN)
    add_footer(s, 5)

    # 6
    s = slides[5]
    add_title(s, "版本迭代路线", "从能降噪到更接近真实听感")
    rows = [
        ["版本阶段", "主要变化", "结论"],
        ["Band-mask TCN", "121K 级别 causal TCN 输出频带 mask", "流程打通，但底噪/人声变小明显"],
        ["空间特征", "加入 IPD、coherence、双麦能量特征", "让模型真正利用双麦信息"],
        ["Tiny DeepFilter", "加入低频复数域多帧滤波", "听感和指标优于纯 mask"],
        ["后处理/GRU/complex", "dehiss、gate、TinyGRU、complex mask", "有局部改善，但不是主解"],
        ["Wind + MWF teacher", "风噪数据 + oracle MWF 蒸馏", "当前推荐主线"],
    ]
    add_table(s, 0.55, 1.4, 12.25, 4.7, rows, [2.15, 5.1, 5.0], 9.3)
    add_footer(s, 6)

    # 7
    s = slides[6]
    add_title(s, "MWF Teacher 的意义", "teacher 只用于训练，不直接部署")
    add_card(s, 0.85, 1.45, 5.5, 2.0, "为什么需要 teacher", ["Oracle local MWF 在风噪评估上达到 +8.82 dB 上限", "说明强双麦前端 / MWF 是有效方向"], GREEN)
    add_card(s, 6.95, 1.45, 5.5, 2.0, "为什么不直接部署", ["Oracle MWF 训练时使用 clean/noise 信息", "真实设备推理时没有 clean，因此只能作为训练老师"], RED)
    add_bullets(s, 1.05, 4.1, 11.2, 1.0, ["当前采用保守蒸馏：target = 0.25 * oracle_mwf_teacher + 0.75 * clean_mic0，避免小模型被过强 teacher 拉向过度滤波。"], 17, INK)
    add_footer(s, 7)

    # 8
    s = slides[7]
    add_title(s, "当前客观指标", "推荐版本：teacher_deepfilter_conservative_on_wind_eval")
    add_metric(s, 0.8, 1.55, 2.7, "Wind noisy", "5.39 dB", "mean noisy SI-SDR", MUTED)
    add_metric(s, 3.45, 1.55, 2.7, "Wind enhanced", "10.19 dB", "mean enhanced SI-SDR", BLUE)
    add_metric(s, 6.1, 1.55, 2.7, "Wind improvement", "+4.80 dB", "当前风噪主结果", GREEN)
    add_metric(s, 8.95, 1.55, 2.7, "Output/Input RMS", "0.764", "人声不过度变小", AMBER)
    rows = [
        ["评估集", "Noisy SI-SDR", "Enhanced SI-SDR", "Improvement"],
        ["Zenodo wind eval", "5.39 dB", "10.19 dB", "+4.80 dB"],
        ["Original DEMAND eval", "4.33 dB", "10.96 dB", "+6.63 dB"],
    ]
    add_table(s, 1.15, 4.15, 10.8, 1.35, rows, [3.3, 2.4, 2.6, 2.5], 11)
    add_footer(s, 8)

    # 9
    s = slides[8]
    add_title(s, "失败样本分析", "现在主要瓶颈变成低频风噪残留")
    rows = [
        ["failure mode", "Wind eval 数量"],
        ["ok", "117"],
        ["low_freq_wind", "35"],
        ["regression", "7"],
        ["residual_noise", "1"],
    ]
    add_table(s, 0.9, 1.5, 4.6, 2.8, rows, [2.6, 2.0], 11)
    add_card(s, 6.25, 1.5, 5.7, 1.35, "关键判断", ["当前模型已经具备降噪能力，但 clean 对比下仍有低频风噪/气流残留。"], AMBER)
    add_card(s, 6.25, 3.2, 5.7, 1.35, "优化方向", ["不能盲目提高全频压噪强度，应针对低频风噪失败样本做定向优化，同时保留语音。"], BLUE)
    add_footer(s, 9)

    # 10
    s = slides[9]
    add_title(s, "网页 Demo 怎么展示", "建议现场展示顺序")
    add_card(s, 0.85, 1.45, 11.65, 0.8, "入口", ["http://127.0.0.1:38180/runs/audio_demo/index.html"], BLUE)
    add_bullets(
        s,
        1.05,
        2.75,
        11.2,
        2.2,
        [
            "先播放 Noisy，让听众感受风噪/背景噪问题。",
            "再播放 teacher_deepfilter_conservative_on_wind_eval 的 Realtime，说明当前模型实际推理效果。",
            "最后播放 Clean，主动说明当前版本仍和 clean 有差距，主要差在低频呼呼声残留。",
            "可补充对比 wind_finetune 和 aggressive teacher，说明为什么选择当前保守 teacher 版本。",
        ],
        15,
        INK,
    )
    add_footer(s, 10)

    # 11
    s = slides[10]
    add_title(s, "阶段性成果", "目前已经形成完整实验闭环")
    add_card(s, 0.75, 1.5, 3.7, 2.1, "算法", ["Tiny band-mask TCN", "Tiny DeepFilter TCN", "coherence/MWF 前端", "MWF teacher 蒸馏"], BLUE)
    add_card(s, 4.85, 1.5, 3.7, 2.1, "数据与评估", ["ARCTIC + DEMAND", "Zenodo Wind Noise", "SI-SDR 评估", "failure analysis"], GREEN)
    add_card(s, 8.95, 1.5, 3.7, 2.1, "工程展示", ["实时增强脚本", "网页试听 Demo", "早期 INT8/C reference", "GitHub 文档同步"], AMBER)
    add_bullets(s, 1.0, 4.55, 11.2, 0.8, ["核心结论：强双麦前端 / MWF teacher 是有效方向，但真实风噪听感还需要低频定向优化和更真实双麦数据。"], 18, INK)
    add_footer(s, 11)

    # 12
    s = slides[11]
    add_title(s, "下一步计划", "围绕当前主线继续收敛")
    add_bullets(
        s,
        1.0,
        1.55,
        11.4,
        4.6,
        [
            "低频风噪定向优化：固定 low_freq_wind 失败样本，加入低频 residual wind loss。",
            "语音保持约束：speech-active RMS / speech-band loss，避免压噪时再次把人声压小。",
            "补充真实双麦数据：佩戴位置、走动风噪、摩擦噪声、麦克风自噪声和真实房间混响。",
            "端侧部署闭环：当前 TinyDeepFilterTCN 的 INT8 导出、fixed-scale Python reference、C reference 对齐和 benchmark。",
            "展示完善：网页 demo 保留固定样本、固定版本对比，主观听感和客观指标一起汇报。",
        ],
        16,
        INK,
    )
    add_footer(s, 12)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
