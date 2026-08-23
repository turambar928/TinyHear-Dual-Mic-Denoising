from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("docs/presentation/tinyhear_project_intro.pptx")


NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
BG = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(203, 213, 225)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
VIOLET = RGBColor(109, 40, 217)


def add_shape(slide, shape_type, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    return shp


def add_text(slide, x, y, w, h, text, size=20, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_para(slide, x, y, w, h, lines, size=18, color=SLATE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_header(slide, title, idx, tag="TinyHear"):
    add_shape(slide, 1, 0, 0, 13.333, 0.18, BLUE)
    add_text(slide, 0.55, 0.42, 9.8, 0.45, title, 27, NAVY, True)
    add_text(slide, 10.65, 0.48, 1.9, 0.25, tag, 11, MUTED, False, PP_ALIGN.RIGHT)
    add_text(slide, 12.35, 6.95, 0.5, 0.25, f"{idx:02d}", 11, MUTED, False, PP_ALIGN.RIGHT)
    add_shape(slide, 1, 0.55, 1.08, 12.1, 0.02, LINE)


def add_card(slide, x, y, w, h, title, lines, accent=BLUE, title_size=17, body_size=14):
    add_shape(slide, 5, x, y, w, h, WHITE, LINE)
    add_shape(slide, 1, x, y, 0.1, h, accent)
    add_text(slide, x + 0.23, y + 0.18, w - 0.45, 0.3, title, title_size, NAVY, True)
    add_para(slide, x + 0.25, y + 0.58, w - 0.45, h - 0.68, lines, body_size, SLATE)


def add_big_metric(slide, x, y, w, h, label, value, note, color=BLUE):
    add_shape(slide, 5, x, y, w, h, WHITE, LINE)
    add_text(slide, x + 0.12, y + 0.15, w - 0.24, 0.25, label, 13, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, x + 0.12, y + 0.47, w - 0.24, 0.55, value, 31, color, True, PP_ALIGN.CENTER)
    add_text(slide, x + 0.12, y + 1.08, w - 0.24, 0.3, note, 12, SLATE, False, PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h, rows, col_widths, size=13):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(226, 232, 240) if r == 0 else WHITE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(size)
                p.font.bold = r == 0
                p.font.color.rgb = NAVY if r == 0 else SLATE
    return tbl


def add_flow(slide, y, items):
    x = 0.55
    w = 1.58
    colors = [BLUE, CYAN, GREEN, AMBER, VIOLET, BLUE, GREEN]
    for i, item in enumerate(items):
        add_shape(slide, 5, x, y, w, 0.8, colors[i % len(colors)])
        add_text(slide, x + 0.08, y + 0.23, w - 0.16, 0.24, item, 12, WHITE, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            add_text(slide, x + w + 0.03, y + 0.25, 0.22, 0.2, ">", 16, MUTED, True, PP_ALIGN.CENTER)
        x += w + 0.24


def add_timeline(slide, x, y, items):
    total_w = 11.9
    step = total_w / (len(items) - 1)
    add_shape(slide, 1, x, y + 0.2, total_w, 0.04, LINE)
    for i, (name, desc, color) in enumerate(items):
        cx = x + i * step
        add_shape(slide, 9, cx - 0.13, y + 0.08, 0.26, 0.26, color)
        add_text(slide, cx - 0.8, y + 0.48, 1.6, 0.28, name, 13, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, cx - 0.95, y + 0.86, 1.9, 0.62, desc, 10.5, SLATE, False, PP_ALIGN.CENTER)


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_shape(s, 1, 0, 0, 13.333, 0.22, BLUE)
    add_shape(s, 1, 0, 6.95, 13.333, 0.22, GREEN)
    add_text(s, 0.72, 0.85, 11.8, 0.65, "微听双麦降噪", 42, NAVY, True)
    add_text(s, 0.76, 1.58, 10.6, 0.36, "TinyHear Dual-Mic Denoising", 22, BLUE, True)
    add_text(s, 0.78, 2.32, 11.2, 0.48, "面向助听器 / 耳戴设备上行链路的端侧 AI 降噪原型", 24, SLATE, True)
    add_card(s, 0.8, 3.35, 3.45, 1.55, "输入输出", ["双麦输入", "单通道增强语音输出"], BLUE, 18, 16)
    add_card(s, 4.85, 3.35, 3.45, 1.55, "实时约束", ["16 kHz", "256 FFT / 64 hop", "4 ms 帧移"], CYAN, 18, 16)
    add_card(s, 8.9, 3.35, 3.45, 1.55, "当前主线", ["137,984 参数", "Tiny DeepFilter + MWF teacher"], GREEN, 18, 15)
    add_text(s, 0.82, 6.38, 11.2, 0.24, "推荐版本：teacher_deepfilter_conservative_on_wind_eval", 14, MUTED, False)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_slide(prs)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "1. 项目背景：为什么做双麦端侧降噪", 2)
    add_card(s, 0.7, 1.55, 3.85, 2.0, "应用场景", ["助听器、耳戴设备、上行通话", "用户直接感知：语音是否清楚、底噪是否刺耳"], BLUE, 18, 15)
    add_card(s, 4.95, 1.55, 3.85, 2.0, "核心噪声", ["环境噪声", "风噪 / 气流呼呼声", "麦克风自噪声和摩擦噪声"], RED, 18, 15)
    add_card(s, 9.2, 1.55, 3.4, 2.0, "端侧约束", ["低延迟", "小模型", "低功耗", "可定点化"], AMBER, 18, 15)
    add_text(s, 1.0, 4.35, 11.2, 0.44, "阶段目标", 25, NAVY, True, PP_ALIGN.CENTER)
    add_para(s, 1.3, 5.0, 10.7, 0.8, ["在 100K-150K 参数级别内，验证双麦空间信息 + 小模型是否能形成实时可展示的语音增强闭环。"], 21, SLATE)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "2. 技术判断：近距双麦不能只看幅度", 3)
    add_card(s, 0.75, 1.45, 5.65, 2.0, "为什么不是单麦降噪", ["双麦有空间信息，可以利用目标语音和噪声方向/相干性的差异", "单麦模型容易把风噪残留成低频呼呼声"], BLUE, 18, 15)
    add_card(s, 6.9, 1.45, 5.65, 2.0, "为什么重点用 IPD / coherence", ["助听器/耳戴设备麦克风距离近，ILD 幅度差通常不明显", "相位差、相干性和局部空间统计更稳定"], GREEN, 18, 15)
    add_shape(s, 5, 1.05, 4.25, 11.25, 1.3, WHITE, LINE)
    add_text(s, 1.35, 4.55, 10.7, 0.38, "当前路线：保留双麦空间前端，同时用轻量神经网络做细化增强", 24, NAVY, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "3. 当前系统整体流程", 4)
    add_flow(s, 1.55, ["双麦输入", "STFT", "IPD/相干性", "MWF 前端", "Tiny DeepFilter", "iSTFT", "网页试听"])
    add_card(s, 0.85, 3.1, 5.65, 2.2, "推理链路", ["输入双通道 mixture wav", "计算双麦空间特征和参考通道", "小模型逐帧输出 gain 与 deep-filter 系数", "实时 overlap-add 得到增强语音"], BLUE, 18, 15)
    add_card(s, 6.95, 3.1, 5.55, 2.2, "评估链路", ["固定风噪 eval 和 DEMAND eval", "输出 SI-SDR、RMS、failure analysis", "网页直接听 Noisy / Realtime / Clean"], GREEN, 18, 15)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "4. 当前推荐模型结构", 5)
    rows = [
        ["模块", "设计"],
        ["前端", "coherence-weighted MWF spatial frontend"],
        ["主干网络", "TinyDeepFilterTCN，channels=96，blocks=8"],
        ["输出", "band gain + low-frequency complex deep-filter coef"],
        ["时频配置", "16 kHz，256 FFT，64 samples hop"],
        ["模型规模", "137,984 参数"],
    ]
    add_table(s, 0.75, 1.45, 11.85, 2.75, rows, [3.0, 8.85], 14)
    add_card(s, 0.95, 4.75, 5.3, 1.3, "band gain 的作用", ["做整体谱幅度抑制，稳定、轻量，便于端侧实现。"], CYAN, 18, 16)
    add_card(s, 7.0, 4.75, 5.3, 1.3, "DeepFilter 的作用", ["低频多帧复数滤波，针对相位和风噪残留问题。"], VIOLET, 18, 16)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "5. 数据集建设：从通用噪声到风噪", 6)
    rows = [
        ["阶段", "Clean speech", "Noise", "作用"],
        ["Baseline", "合成语音", "合成噪声", "快速打通训练/推理"],
        ["小规模验证", "YESNO", "DEMAND", "sanity check"],
        ["主基线", "CMU ARCTIC", "DEMAND", "环境噪声基线"],
        ["当前主线", "CMU ARCTIC", "DEMAND + Zenodo Wind", "风噪重点优化"],
    ]
    add_table(s, 0.65, 1.35, 12.0, 3.15, rows, [2.2, 2.65, 3.0, 4.15], 12.5)
    add_card(s, 0.95, 5.05, 5.4, 1.0, "训练目录", ["data/arctic_wind_demand_flat/"], BLUE, 17, 15)
    add_card(s, 6.95, 5.05, 5.4, 1.0, "固定评估集", ["data/arctic_wind_eval/val；data/arctic_demand_eval/val"], GREEN, 17, 15)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "6. 版本迭代总览", 7)
    add_timeline(
        s,
        0.8,
        1.55,
        [
            ("Baseline", "band-mask TCN\n能降噪但底噪大", BLUE),
            ("Spatial", "IPD / coherence\n利用双麦信息", CYAN),
            ("DeepFilter", "多帧复数滤波\n缓解相位问题", GREEN),
            ("Postfilter", "dehiss / gate\n只能局部改善", AMBER),
            ("Wind data", "加入真实风噪\n听感改善", VIOLET),
            ("MWF teacher", "当前推荐主线\n风噪更稳", RED),
        ],
    )
    add_card(s, 0.9, 4.25, 11.55, 1.5, "迭代结论", ["反复实验后，单靠后处理或换一个小模型不能根治呼呼声；更有效的方向是“更强双麦前端 + DeepFilter + MWF teacher 蒸馏”。"], BLUE, 20, 17)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "7. 为什么引入 MWF Teacher", 8)
    add_big_metric(s, 0.85, 1.55, 3.2, 1.55, "Oracle local MWF", "+8.82 dB", "Wind eval 上限验证", GREEN)
    add_card(s, 4.55, 1.55, 3.9, 1.55, "teacher 是什么", ["训练阶段用 clean/noisy 构造的强滤波器", "用于告诉小模型更合理的滤波方向"], BLUE, 18, 15)
    add_card(s, 8.9, 1.55, 3.6, 1.55, "为什么不能部署", ["真实推理没有 clean 信息", "最终部署仍是 137,984 参数小模型"], RED, 18, 15)
    add_text(s, 1.0, 4.05, 11.3, 0.4, "当前采用保守蒸馏", 24, NAVY, True, PP_ALIGN.CENTER)
    add_shape(s, 5, 2.2, 4.75, 8.9, 0.75, WHITE, LINE)
    add_text(s, 2.45, 4.98, 8.4, 0.28, "target = 0.25 * oracle_mwf_teacher + 0.75 * clean_mic0", 20, BLUE, True, PP_ALIGN.CENTER)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "8. 训练目标与 Loss 设计", 9)
    add_card(s, 0.75, 1.45, 3.75, 2.0, "重建质量", ["waveform L1", "SI-SDR loss", "multi-resolution / log STFT magnitude"], BLUE, 18, 15)
    add_card(s, 4.8, 1.45, 3.75, 2.0, "残留噪声控制", ["residual noise loss", "silence floor loss", "低能量区域底噪约束"], GREEN, 18, 15)
    add_card(s, 8.85, 1.45, 3.75, 2.0, "稳定性", ["coef energy regularization", "保守 teacher blend", "避免过度滤波"], AMBER, 18, 15)
    add_para(s, 1.0, 4.55, 11.2, 1.0, ["设计原则：不能只优化 MSE 或单一 SI-SDR，因为听感中的风噪、气流声和人声变小往往不是一个指标能完整反映。"], 21, SLATE)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "9. 当前客观结果", 10)
    add_big_metric(s, 0.85, 1.45, 2.85, 1.55, "Wind noisy", "5.39", "SI-SDR dB", MUTED)
    add_big_metric(s, 3.95, 1.45, 2.85, 1.55, "Wind enhanced", "10.19", "SI-SDR dB", BLUE)
    add_big_metric(s, 7.05, 1.45, 2.85, 1.55, "Wind gain", "+4.80", "SI-SDR improvement", GREEN)
    add_big_metric(s, 10.15, 1.45, 2.4, 1.55, "RMS ratio", "0.764", "输出响度保持", AMBER)
    rows = [
        ["评估集", "Noisy", "Enhanced", "Improvement"],
        ["Zenodo wind eval", "5.39 dB", "10.19 dB", "+4.80 dB"],
        ["Original DEMAND eval", "4.33 dB", "10.96 dB", "+6.63 dB"],
    ]
    add_table(s, 1.15, 4.1, 11.0, 1.35, rows, [3.4, 2.4, 2.6, 2.6], 13.5)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "10. 关键版本对比", 11)
    rows = [
        ["版本", "Wind eval", "DEMAND eval", "结论"],
        ["wind fine-tune", "+4.35 dB", "+7.06 dB", "通用噪声更稳，风噪仍有残留"],
        ["aggressive teacher", "+3.74 dB", "+5.21 dB", "teacher 太强，小模型学不稳"],
        ["conservative teacher", "+4.80 dB", "+6.63 dB", "当前推荐，风噪更优"],
    ]
    add_table(s, 0.75, 1.5, 11.8, 2.3, rows, [3.15, 2.1, 2.15, 4.4], 13)
    add_card(s, 1.05, 4.55, 11.25, 1.2, "选择当前版本的原因", ["它不是所有指标都最高，而是在风噪听感、人声保持和模型稳定性之间更均衡。"], GREEN, 20, 18)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "11. 失败样本分析", 12)
    rows = [
        ["Failure mode", "Wind eval 数量"],
        ["ok", "117"],
        ["low_freq_wind", "35"],
        ["regression", "7"],
        ["residual_noise", "1"],
    ]
    add_table(s, 0.95, 1.5, 4.8, 3.1, rows, [2.9, 1.9], 14)
    add_card(s, 6.45, 1.55, 5.65, 1.45, "诊断结论", ["当前主要问题不是完全没有降噪能力，而是低频风噪/气流残留。"], AMBER, 20, 17)
    add_card(s, 6.45, 3.45, 5.65, 1.45, "后续方向", ["围绕 low_freq_wind 失败样本做定向 loss 和真实数据补充。"], BLUE, 20, 17)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "12. 网页 Demo 展示方式", 13)
    add_card(s, 0.85, 1.4, 11.75, 0.95, "网页入口", ["http://127.0.0.1:38180/runs/audio_demo/index.html"], BLUE, 19, 17)
    add_card(s, 0.95, 2.75, 3.55, 2.25, "第一步", ["播放 Noisy", "让听众先听到原始风噪和底噪问题"], RED, 19, 16)
    add_card(s, 4.9, 2.75, 3.55, 2.25, "第二步", ["播放 Realtime", "选择 teacher_deepfilter_conservative_on_wind_eval"], GREEN, 19, 16)
    add_card(s, 8.85, 2.75, 3.55, 2.25, "第三步", ["播放 Clean", "说明当前差距主要在低频呼呼声残留"], BLUE, 19, 16)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "13. 已完成工程闭环", 14)
    add_card(s, 0.75, 1.45, 3.75, 2.0, "训练", ["数据准备", "动态双麦混音", "多版本模型训练", "teacher 蒸馏"], BLUE, 18, 15)
    add_card(s, 4.8, 1.45, 3.75, 2.0, "评估", ["SI-SDR metrics", "实时增强脚本", "failure analysis", "固定样本对比"], GREEN, 18, 15)
    add_card(s, 8.85, 1.45, 3.75, 2.0, "展示/部署", ["网页试听 demo", "GitHub 文档", "早期 INT8 导出", "C reference 链路"], AMBER, 18, 15)
    add_para(s, 1.0, 4.55, 11.2, 0.8, ["目前项目已经不是单个模型实验，而是包含数据、训练、评估、试听和工程化验证的完整原型系统。"], 21, SLATE)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "14. 当前不足", 15)
    add_card(s, 0.8, 1.4, 5.65, 1.35, "听感差距", ["相比 clean，部分样本仍有低频风噪和气流呼呼声。"], RED, 20, 17)
    add_card(s, 6.85, 1.4, 5.65, 1.35, "数据差距", ["公开数据集合成与真实助听器佩戴、真实风噪仍有 domain gap。"], AMBER, 20, 17)
    add_card(s, 0.8, 3.3, 5.65, 1.35, "指标差距", ["SI-SDR 与主观听感不完全一致，需要固定听感评估集。"], BLUE, 20, 17)
    add_card(s, 6.85, 3.3, 5.65, 1.35, "部署差距", ["当前推荐 DeepFilter 版本还需要完成 INT8/C reference 对齐。"], GREEN, 20, 17)

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_header(s, "15. 下一步计划", 16)
    add_para(
        s,
        0.95,
        1.45,
        11.7,
        4.6,
        [
            "1. 低频风噪定向优化：固定 low_freq_wind 失败样本，增加低频 residual wind loss。",
            "2. 语音保持约束：speech-active RMS / speech-band loss，避免继续压噪时人声变小。",
            "3. 补充真实双麦数据：佩戴位置、走动风噪、摩擦噪声、麦克风自噪声和房间混响。",
            "4. 端侧部署闭环：TinyDeepFilterTCN 的 INT8 导出、fixed-scale Python reference、C reference 对齐和 benchmark。",
            "5. 展示完善：固定 demo 样本和版本对比，主观听感 + 客观指标同时汇报。",
        ],
        18,
        SLATE,
    )

    s = prs.slides.add_slide(blank)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    add_shape(s, 1, 0, 0, 13.333, 0.22, BLUE)
    add_text(s, 0.75, 0.9, 11.9, 0.55, "阶段性结论", 36, NAVY, True)
    add_shape(s, 5, 0.95, 1.95, 11.45, 2.1, WHITE, LINE)
    add_text(s, 1.35, 2.35, 10.65, 0.35, "项目已经验证：强双麦前端 / MWF teacher 是有效方向", 25, GREEN, True, PP_ALIGN.CENTER)
    add_para(
        s,
        1.25,
        4.45,
        10.9,
        1.2,
        [
            "当前推荐版本在风噪 eval 上达到 +4.80 dB SI-SDR improvement，在 DEMAND eval 上达到 +6.63 dB。",
            "下一阶段重点不是盲目换模型，而是围绕低频风噪残留、真实双麦数据和端侧部署闭环继续收敛。",
        ],
        21,
        SLATE,
    )
    add_text(s, 12.2, 6.95, 0.55, 0.25, "17", 11, MUTED, False, PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
